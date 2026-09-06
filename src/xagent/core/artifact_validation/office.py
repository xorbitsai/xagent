"""OOXML package preflight, followed by the actual document readers."""

from io import BytesIO
from pathlib import Path, PurePosixPath
from zipfile import BadZipFile, ZipFile

from .models import ArtifactContent, InvalidArtifact, UncheckedArtifact


def check_package(content: ArtifactContent) -> None:
    from xml.etree.ElementTree import ParseError

    from defusedxml.common import DefusedXmlException
    from defusedxml.ElementTree import fromstring

    try:
        with ZipFile(BytesIO(content.data)) as archive:
            entries = archive.infolist()
            if (
                len(entries) > content.limits.max_entries
                or sum(e.file_size for e in entries) > content.limits.max_expanded_bytes
            ):
                raise UncheckedArtifact("Office package exceeds the expansion budget.")
            names = [e.filename for e in entries]
            if len(names) != len(set(names)) or any(
                PurePosixPath(n).is_absolute() or ".." in PurePosixPath(n).parts
                for n in names
            ):
                raise InvalidArtifact("Office package contains ambiguous member paths.")
            # Main-part locations come from content types and relationships,
            # not conventional filenames. The format reader resolves them
            # after this package-wide safety preflight has passed.
            required = {"[Content_Types].xml", "_rels/.rels"}
            if not required.issubset(names):
                raise InvalidArtifact(
                    "Office package is missing required document parts."
                )
            expanded = 0
            for entry in entries:
                if entry.flag_bits & 1:
                    raise UncheckedArtifact(
                        "Encrypted Office packages cannot be checked."
                    )
                # Read all members to exercise CRC/truncation checks, without
                # extracting files or trusting compressed-size claims alone.
                with archive.open(entry) as stream:
                    data = stream.read(content.limits.max_expanded_bytes - expanded + 1)
                expanded += len(data)
                if expanded > content.limits.max_expanded_bytes:
                    raise UncheckedArtifact(
                        "Office package exceeds the expansion budget."
                    )
                if entry.filename.endswith((".xml", ".rels")):
                    fromstring(data)
    except (BadZipFile, ParseError, DefusedXmlException, EOFError, OSError) as exc:
        raise InvalidArtifact(
            "Office package is corrupt or contains unsafe XML."
        ) from exc


def check_office_document(content: ArtifactContent) -> None:
    suffix = Path(content.filename).suffix.lower()
    try:
        if suffix == ".xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(
                BytesIO(content.data), read_only=True, data_only=False, keep_links=False
            )
            try:
                cells = 0
                for sheet in workbook:
                    # Ignore stale dimension metadata; walk actual cells. Empty
                    # workbooks/templates are legal, not business failures.
                    sheet.reset_dimensions()
                    for row in sheet.iter_rows():
                        cells += len(row)
                        if cells > content.limits.max_units:
                            raise UncheckedArtifact("Workbook exceeds the cell budget.")
            finally:
                workbook.close()
        elif suffix == ".docx":
            from docx import Document

            document = Document(BytesIO(content.data))
            _ = document.paragraphs, document.tables
        else:
            from pptx import Presentation

            presentation = Presentation(BytesIO(content.data))
            for slide in presentation.slides:
                _ = slide.shapes
    except (ValueError, KeyError, IndexError, TypeError, OSError, BadZipFile) as exc:
        raise InvalidArtifact(
            "Office document cannot be opened by its format reader."
        ) from exc
