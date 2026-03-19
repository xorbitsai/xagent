"""File upload API route handlers"""

import asyncio
import logging
import re
import ipaddress
import time
from urllib.parse import urlparse
from pathlib import Path
from typing import Any, Dict

from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, Response, StreamingResponse
from pptx import Presentation
from sqlalchemy.orm import Session
from starlette.concurrency import run_in_threadpool

from ...core.tools.adapters.vibe.file_tool import read_file
from ..auth_dependencies import get_current_user
from ..config import (
    MAX_FILE_SIZE,
    UPLOADS_DIR,
    get_file_info,
    get_file_url,
    get_upload_path,
    is_allowed_file,
)
from ..models.database import get_db
from ..models.user import User

logger = logging.getLogger(__name__)

# Max size (bytes) to serve in-memory for preview to avoid Content-Length mismatch
# (e.g. file changed between stat and read, or NFS/volume sync issues)
PREVIEW_IN_MEMORY_MAX_BYTES = 10 * 1024 * 1024  # 10MB
# Images always served in-memory (up to this cap) to avoid intermittent ERR_CONTENT_LENGTH_MISMATCH
PREVIEW_IMAGE_MAX_BYTES = 50 * 1024 * 1024  # 50MB


def _is_file_being_written(file_path: Path) -> bool:
    """Check if a file is potentially being written by looking for temp files."""
    # Check for common temp file patterns
    temp_patterns = [
        file_path.with_suffix(file_path.suffix + ".tmp"),
        file_path.with_suffix(file_path.suffix + ".temp"),
        file_path.parent / (file_path.name + ".tmp"),
        file_path.parent / (file_path.name + ".temp"),
    ]
    for temp_file in temp_patterns:
        if temp_file.exists():
            return True
    return False


def _read_file_with_retry(file_path: Path, max_retries: int = 5, delay: float = 0.1) -> bytes:
    """
    Read file contents with retry logic to handle race conditions.

    This helps avoid ERR_CONTENT_LENGTH_MISMATCH when:
    - File is being written while we're reading
    - NFS/volume sync delays in Docker environments
    - File is being moved/renamed

    Args:
        file_path: Path to the file to read
        max_retries: Maximum number of retry attempts
        delay: Initial delay between retries (doubles each retry)

    Returns:
        File contents as bytes

    Raises:
        IOError: If file cannot be read after all retries
    """
    last_error = None
    for attempt in range(max_retries):
        try:
            # First check if file exists
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")

            # Check if file might still be being written (temp file exists)
            if _is_file_being_written(file_path):
                raise IOError(f"File appears to be in progress (temp file detected): {file_path}")

            # Get file size before reading
            size_before = file_path.stat().st_size

            # Read file into memory completely before returning
            # This ensures we get a consistent snapshot
            content = file_path.read_bytes()

            # Verify the content matches expected size (if file hasn't changed)
            size_after = file_path.stat().st_size
            if len(content) != size_before or len(content) != size_after:
                # File changed during read, retry
                raise IOError(
                    f"File size mismatch: read {len(content)} bytes, "
                    f"size_before={size_before}, size_after={size_after}"
                )

            # Additional validation: check if content is not empty for images
            if len(content) == 0:
                raise IOError(f"File is empty: {file_path}")

            logger.debug(f"Successfully read file {file_path} ({len(content)} bytes) on attempt {attempt + 1}")
            return content

        except (IOError, OSError) as e:
            last_error = e
            if attempt < max_retries - 1:
                wait_time = delay * (2 ** attempt)  # Exponential backoff
                logger.warning(f"File read attempt {attempt + 1} failed for {file_path}: {e}. Retrying in {wait_time}s...")
                time.sleep(wait_time)
            else:
                logger.error(f"File read failed after {max_retries} attempts for {file_path}: {e}")

    raise last_error or IOError(f"Failed to read file: {file_path}")

# CORS headers for file endpoints (avoid CORS / ERR_CONTENT_LENGTH_MISMATCH when loading in <img>)
FILE_RESPONSE_CORS_HEADERS = {"Access-Control-Allow-Origin": "*"}

# Create router
file_router = APIRouter(prefix="/api/files", tags=["files"])

PROXY_IMAGE_MAX_BYTES = 20 * 1024 * 1024  # 20MB (avoid huge downloads via proxy)


def _extract_task_id_from_web_task_path(file_path: str) -> int | None:
    """
    Extract numeric task_id from paths like:
    - web_task_47/output/foo.jpg
    - /web_task_47/output/foo.jpg
    - web_task_47%2Foutput%2Ffoo.jpg (if passed pre-encoded)
    """
    if not file_path:
        return None
    m = re.search(r"web_task_(\d+)", file_path)
    if not m:
        return None
    try:
        return int(m.group(1))
    except ValueError:
        return None


def _is_private_host(hostname: str) -> bool:
    """
    Best-effort SSRF guard for obvious private targets.
    Note: does not perform DNS resolution (safe/fast); blocks literal IPs and common local names.
    """
    if not hostname:
        return True
    h = hostname.strip().lower()
    if h in {"localhost"}:
        return True
    # literal IP?
    try:
        ip = ipaddress.ip_address(h)
        return (
            ip.is_private
            or ip.is_loopback
            or ip.is_link_local
            or ip.is_multicast
            or ip.is_unspecified
            or ip.is_reserved
        )
    except ValueError:
        # not an IP, allow (DNS resolution could still map to private, but we avoid doing DNS here)
        return False


@file_router.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    task_type: str = Form(...),
    message: str = Form(""),
    task_id: str = Form(None),
    folder: str = Form(None),
    user: User = Depends(get_current_user),
) -> Dict:
    """
    Upload a single file and optionally create a task (backward compatibility)

    Args:
        file: Uploaded file
        task_type: Type of task (e.g., "general")
        message: Optional message to include with the task

    Returns:
        Upload result with file info and optional task ID
    """
    try:
        uploaded_files = []

        # Process the file
        if not file.filename or not file.filename.strip():
            raise HTTPException(status_code=422, detail="No filename provided")

        # Check file extension
        if not is_allowed_file(file.filename, task_type):
            raise HTTPException(
                status_code=500,
                detail=f"File type {Path(file.filename).suffix.lower()} not supported for task type {task_type}",
            )

        # Check file size
        content = await file.read()
        if len(content) > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=500,
                detail=f"File size exceeds maximum limit of {MAX_FILE_SIZE // (1024 * 1024)}MB",
            )

        # Get upload path with user isolation
        file_path = get_upload_path(file.filename, task_id, folder, int(user.id))

        # Save uploaded file
        with open(file_path, "wb") as buffer:
            buffer.write(content)

        logger.info(f"File uploaded: {file.filename} -> {file_path} (user: {user.id})")

        # Read file content for processing
        try:
            file_content = read_file(str(file_path))

            uploaded_files.append(
                {
                    "filename": file.filename,
                    "file_path": str(file_path),
                    "file_url": get_file_url(
                        file.filename, task_id, folder, int(user.id)
                    ),
                    "file_size": len(content),
                    "content_preview": file_content[:500] + "..."
                    if len(file_content) > 500
                    else file_content,
                }
            )

        except (ValueError, KeyError, TypeError) as e:
            # Data format error
            logger.error(f"Data format error processing file {file.filename}: {e}")
            uploaded_files.append(
                {
                    "filename": file.filename,
                    "file_path": str(file_path),
                    "file_url": get_file_url(
                        file.filename, task_id, folder, int(user.id)
                    ),
                    "file_size": len(content),
                    "error": f"Data format error: {str(e)}",
                }
            )
        except (PermissionError, OSError) as e:
            # File system permission error
            logger.error(f"File system error processing file {file.filename}: {e}")
            uploaded_files.append(
                {
                    "filename": file.filename,
                    "file_path": str(file_path),
                    "file_url": get_file_url(
                        file.filename, task_id, folder, int(user.id)
                    ),
                    "file_size": len(content),
                    "error": f"File system error: {str(e)}",
                }
            )
        except Exception as e:
            # Other errors, re-raise
            logger.error(f"Unexpected error processing file {file.filename}: {e}")
            raise

        # Determine overall success
        all_successful = all("error" not in file_info for file_info in uploaded_files)

        # Return single file format
        file_info = uploaded_files[0]
        return {
            "success": all_successful,
            "filename": file_info["filename"],
            "file_path": file_info["file_path"],
            "file_url": file_info["file_url"],
            "file_size": file_info["file_size"],
            "task_type": task_type,
            "content_preview": file_info.get("content_preview", ""),
            "error": file_info.get("error"),
            "message": f"Successfully uploaded {file_info['filename']}"
            if all_successful
            else f"Failed to process {file_info['filename']}",
        }

    except HTTPException:
        # Re-raise HTTP exceptions (like 422 validation errors)
        raise
    except (ValueError, KeyError, TypeError) as e:
        # Data format error
        logger.error(f"Data format error in file upload: {e}")
        raise HTTPException(status_code=400, detail=f"Data format error: {str(e)}")
    except (PermissionError, OSError) as e:
        # File system permission error
        logger.error(f"File system error in file upload: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except Exception as e:
        # Other errors, re-raise
        logger.error(f"Unexpected error in file upload: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.post("/upload-multiple")
async def upload_multiple_files(
    files: list[UploadFile] = File(...),
    task_type: str = Form(...),
    message: str = Form(""),
    task_id: str = Form(None),
    folder: str = Form(None),
    user: User = Depends(get_current_user),
) -> Dict:
    """
    Upload multiple files and optionally create a task (new format)

    Args:
        files: List of uploaded files
        task_type: Type of task (e.g., "general")
        message: Optional message to include with the task
        task_id: Optional task ID to organize files in task-specific folder
        folder: Optional folder name within task directory (e.g., "input", "output")

    Returns:
        Upload result with file info and optional task ID
    """
    try:
        uploaded_files = []

        # Process each file
        for file in files:
            # Validate file type
            if not file.filename or not file.filename.strip():
                raise HTTPException(status_code=422, detail="No filename provided")

            # Check file extension
            if not is_allowed_file(file.filename, task_type):
                raise HTTPException(
                    status_code=500,
                    detail=f"File type {Path(file.filename).suffix.lower()} not supported for task type {task_type}",
                )

            # Check file size
            content = await file.read()
            if len(content) > MAX_FILE_SIZE:
                raise HTTPException(
                    status_code=500,
                    detail=f"File size exceeds maximum limit of {MAX_FILE_SIZE // (1024 * 1024)}MB",
                )

            # Get upload path with user isolation
            file_path = get_upload_path(file.filename, task_id, folder, int(user.id))

            # Save uploaded file
            with open(file_path, "wb") as buffer:
                buffer.write(content)

            logger.info(
                f"File uploaded: {file.filename} -> {file_path} (user: {user.id})"
            )

            # Read file content for processing
            try:
                file_content = read_file(str(file_path))

                uploaded_files.append(
                    {
                        "filename": file.filename,
                        "file_path": str(file_path),
                        "file_url": get_file_url(
                            file.filename, task_id, folder, int(user.id)
                        ),
                        "file_size": len(content),
                        "content_preview": file_content[:500] + "..."
                        if len(file_content) > 500
                        else file_content,
                    }
                )

            except (ValueError, KeyError, TypeError) as e:
                # Data format error
                logger.error(f"Data format error processing file {file.filename}: {e}")
                uploaded_files.append(
                    {
                        "filename": file.filename,
                        "file_path": str(file_path),
                        "file_url": get_file_url(
                            file.filename, task_id, folder, int(user.id)
                        ),
                        "file_size": len(content),
                        "error": f"Data format error: {str(e)}",
                    }
                )
            except (PermissionError, OSError) as e:
                # File system permission error
                logger.error(f"File system error processing file {file.filename}: {e}")
                uploaded_files.append(
                    {
                        "filename": file.filename,
                        "file_path": str(file_path),
                        "file_url": get_file_url(
                            file.filename, task_id, folder, int(user.id)
                        ),
                        "file_size": len(content),
                        "error": f"File system error: {str(e)}",
                    }
                )
            except Exception as e:
                # Other errors, re-raise
                logger.error(f"Unexpected error processing file {file.filename}: {e}")
                raise

        # Determine overall success
        all_successful = all("error" not in file_info for file_info in uploaded_files)

        # Return new format for multiple files
        return {
            "success": all_successful,
            "files": uploaded_files,
            "total_files": len(uploaded_files),
            "task_type": task_type,
            "message": f"Successfully uploaded {len(uploaded_files)} files"
            if all_successful
            else "Some files had processing errors",
        }

    except HTTPException:
        # Re-raise HTTP exceptions (like 422 validation errors)
        raise
    except (ValueError, KeyError, TypeError) as e:
        # Data format error
        logger.error(f"Data format error in multiple file upload: {e}")
        raise HTTPException(status_code=400, detail=f"Data format error: {str(e)}")
    except (PermissionError, OSError) as e:
        # File system permission error
        logger.error(f"File system error in multiple file upload: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except Exception as e:
        # Other errors, re-raise
        logger.error(f"Unexpected error in multiple file upload: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.get("/list")
async def list_files(user: User = Depends(get_current_user)) -> Dict:
    """List user's uploaded files"""
    try:
        # Get user-specific directory, or all directories for admin
        if user.is_admin:
            # Admin can see all files - scan all user directories
            scan_dirs = [
                d
                for d in UPLOADS_DIR.iterdir()
                if d.is_dir() and d.name.startswith("user_")
            ]
        else:
            # Regular users can only see their own files
            user_dir = UPLOADS_DIR / f"user_{user.id}"
            scan_dirs = [user_dir] if user_dir.exists() else []

        files = []

        def scan_directory(
            directory: Path, relative_path: str = "", user_prefix: str = ""
        ) -> None:
            """Recursively scan directory for files"""
            for item in directory.iterdir():
                if item.name.startswith("."):
                    continue

                if item.is_file():
                    file_info = get_file_info(str(item))
                    if file_info:
                        # Add relative path for display
                        if relative_path:
                            file_info["relative_path"] = f"{relative_path}/{item.name}"
                        else:
                            file_info["relative_path"] = item.name

                        # Parse relative_path to extract task_id and folder
                        rel_path = file_info.get("relative_path", item.name)
                        path_parts = rel_path.split("/")

                        if len(path_parts) >= 3 and path_parts[0].startswith(
                            "web_task_"
                        ):
                            # Format: web_task_X/folder/filename
                            task_id = path_parts[0].replace("web_task_", "")
                            folder = "/".join(path_parts[1:-1])
                            filename = path_parts[-1]

                            # Extract user_id from user_prefix for admin
                            file_user_id = (
                                int(user_prefix.replace("user_", ""))
                                if user_prefix
                                else int(user.id)
                            )
                            file_info["file_url"] = get_file_url(
                                filename,
                                task_id=task_id,
                                folder=folder,
                                user_id=file_user_id,
                            )
                        else:
                            # Fallback to simple filename
                            file_user_id = (
                                int(user_prefix.replace("user_", ""))
                                if user_prefix
                                else int(user.id)
                            )
                            file_info["file_url"] = get_file_url(
                                file_info["filename"], user_id=file_user_id
                            )

                        # Add user info for admin to identify file ownership
                        if user.is_admin and user_prefix:
                            file_info["user_id"] = int(user_prefix.replace("user_", ""))

                        files.append(file_info)
                elif item.is_dir():
                    # Recursively scan subdirectories
                    new_relative_path = (
                        f"{relative_path}/{item.name}" if relative_path else item.name
                    )
                    scan_directory(item, new_relative_path, user_prefix)

        # Scan all user directories
        for scan_dir in scan_dirs:
            user_prefix = scan_dir.name if scan_dir.name.startswith("user_") else ""
            scan_directory(scan_dir, "", user_prefix)

        return {"files": files, "total_count": len(files)}

    except (PermissionError, OSError) as e:
        # File system permission error
        logger.error(f"File system error listing files: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except (ValueError, KeyError, TypeError) as e:
        # Data format error
        logger.error(f"Data format error listing files: {e}")
        raise HTTPException(status_code=400, detail=f"Data format error: {str(e)}")
    except Exception as e:
        # Other errors, re-raise
        logger.error(f"Unexpected error listing files: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.get("/download/{file_path:path}", response_model=None)
async def download_file(
    file_path: str,
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Any:
    """Download uploaded file"""
    try:
        logger.info(f"Download request for file_path: {file_path}")

        # Check if this is a web_task file and handle admin access
        target_user_dir = None
        if file_path.startswith("web_task_"):
            # Extract task ID from path like "web_task_78/output/file.html"
            try:
                task_id = int(file_path.split("_")[2].split("/")[0])
                from ..models.task import Task

                task = db.query(Task).filter(Task.id == task_id).first()

                if task:
                    # Admin can access any task, regular users can only access their own tasks
                    if not user.is_admin and task.user_id != int(user.id):
                        logger.warning(
                            f"User {user.id} attempted to access task {task_id} belonging to user {task.user_id}"
                        )
                        raise HTTPException(status_code=403, detail="Access denied")

                    # Use the task owner's directory for file access
                    target_user_dir = UPLOADS_DIR / f"user_{task.user_id}"
                    logger.info(
                        f"Accessing task {task_id} files from user directory: user_{task.user_id}"
                    )
                else:
                    logger.warning(f"Task {task_id} not found in database")
                    raise HTTPException(status_code=404, detail="Task not found")
            except (ValueError, IndexError) as e:
                logger.warning(f"Invalid web_task path format: {file_path}, error: {e}")
                # Fall back to normal processing if path format is invalid
                pass

        # Get user-specific directory
        # For admin accessing other users' task files, use the target user's directory
        if target_user_dir:
            user_dir = target_user_dir
        else:
            user_dir = UPLOADS_DIR / f"user_{user.id}"

        # Define recursive search function once
        def find_file_recursively(directory: Path, filename: str) -> Path | None:
            """Recursively find file in directory"""
            for item in directory.iterdir():
                if item.is_file() and item.name == filename:
                    return item
                elif item.is_dir():
                    result = find_file_recursively(item, filename)
                    if result:
                        return result
            return None

        if "/" not in file_path:
            # Simple filename - search in user directory first, then fallback to global
            found_path = None

            # Search in user directory first
            if user_dir.exists():
                found_path = find_file_recursively(user_dir, file_path)
                if found_path:
                    full_path = found_path
                    logger.info(f"Found file in user directory: {full_path}")

            # If not found in user directory, search in global directory for legacy files
            if not found_path:
                if UPLOADS_DIR.exists():
                    found_path = find_file_recursively(UPLOADS_DIR, file_path)
                    if found_path:
                        full_path = found_path
                        logger.info(
                            f"Found file in global directory (legacy): {full_path}"
                        )

            if not found_path:
                logger.warning(f"File not found for download: {file_path}")
                raise HTTPException(status_code=404, detail="File not found")
        else:
            # Relative path - check if it starts with user_id prefix
            if file_path.startswith(f"user_{user.id}/"):
                # Path already includes user directory, use it relative to UPLOADS_DIR
                full_path = UPLOADS_DIR / file_path
            elif target_user_dir:
                # For admin accessing other users' task files, path is relative to target user directory
                full_path = target_user_dir / file_path
            else:
                # Path is relative to user directory
                full_path = user_dir / file_path

            if not full_path.exists():
                # Fallback: search recursively for the filename with path preference
                target_filename = file_path.split("/")[-1]

                # Enhanced search that prefers matching path structure
                def find_file_with_path_preference(
                    directory: Path, filename: str, original_path: str
                ) -> Path | None:
                    """Find file recursively, preferring paths that match the original structure"""
                    best_match = None
                    best_score = 0

                    def score_match(candidate_path: Path, target_path: str) -> int:
                        """Score how well a candidate path matches the target structure"""
                        # Get relative path from user directory
                        try:
                            relative_path = candidate_path.relative_to(user_dir)
                        except ValueError:
                            # Fallback to using absolute path parts
                            relative_path = candidate_path

                        candidate_parts = relative_path.parts
                        target_parts = target_path.split("/")

                        score = 0
                        # Check if path parts match
                        for i, target_part in enumerate(
                            target_parts[:-1]
                        ):  # Exclude filename
                            if (
                                i < len(candidate_parts) - 1
                            ):  # Exclude filename from candidate
                                candidate_part = candidate_parts[i]
                                if target_part == candidate_part:
                                    score += 10
                                elif (
                                    target_part.replace("web_task_", "task_")
                                    == candidate_part
                                ):
                                    score += (
                                        5  # Partial match for task prefix differences
                                    )
                                elif (
                                    candidate_part.replace("web_task_", "task_")
                                    == target_part
                                ):
                                    score += 5  # Reverse match
                        return score

                    for item in directory.iterdir():
                        if item.is_file() and item.name == filename:
                            # Score this match
                            current_score = score_match(item, original_path)
                            if current_score > best_score:
                                best_match = item
                                best_score = current_score
                                logger.debug(
                                    f"Found better match with score {current_score}: {item}"
                                )
                        elif item.is_dir():
                            result = find_file_with_path_preference(
                                item, filename, original_path
                            )
                            if result:
                                current_score = score_match(result, original_path)
                                if current_score > best_score:
                                    best_match = result
                                    best_score = current_score

                    return best_match

                found_path = find_file_with_path_preference(
                    user_dir, target_filename, file_path
                )
                if found_path:
                    full_path = found_path
                    logger.info(f"Found file by enhanced recursive search: {full_path}")
                else:
                    logger.warning(f"File not found for download: {file_path}")
                    raise HTTPException(status_code=404, detail="File not found")
            else:
                logger.info(f"Found file by direct path: {full_path}")

        # Security check: ensure the path is within allowed directories
        try:
            resolved_path = full_path.resolve()
            resolved_user_dir = user_dir.resolve()
            resolved_global_dir = UPLOADS_DIR.resolve()

            # Allow both user directory and global directory for legacy files
            try:
                resolved_path.relative_to(resolved_user_dir)
            except ValueError:
                try:
                    resolved_path.relative_to(resolved_global_dir)
                except ValueError:
                    logger.warning(f"Security check failed for path: {file_path}")
                    raise HTTPException(status_code=403, detail="Access denied")
        except ValueError:
            logger.warning(f"Security check failed for path: {file_path}")
            raise HTTPException(status_code=403, detail="Access denied")

        if not full_path.exists():
            logger.warning(f"File not found for download: {full_path}")
            raise HTTPException(status_code=404, detail="File not found")

        filename = full_path.name

        # Handle PPTX files - convert to PDF using LibreOffice for preview
        if filename.endswith(".pptx"):
            logger.info(f"Converting PPTX to PDF for download: {full_path}")
            import tempfile

            try:
                # Create a temporary directory for the conversion
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Use LibreOffice to convert PPTX to PDF (async)
                    proc = await asyncio.create_subprocess_exec(
                        "soffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        temp_dir,
                        str(full_path),
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                    )

                    try:
                        stdout, stderr = await asyncio.wait_for(
                            proc.communicate(), timeout=30
                        )
                    except asyncio.TimeoutError:
                        proc.kill()
                        await proc.wait()
                        logger.error(
                            f"LibreOffice conversion timed out for PPTX: {full_path}"
                        )
                        raise

                    if proc.returncode == 0:
                        # Find the generated PDF file
                        pdf_files = list(Path(temp_dir).glob("*.pdf"))
                        if pdf_files:
                            pdf_path = pdf_files[0]
                            logger.info(
                                f"Successfully converted PPTX to PDF: {pdf_path}"
                            )

                            # Read PDF into memory before temp dir is cleaned up
                            pdf_content = pdf_path.read_bytes()

                            # Return streaming response from memory (Content-Length avoids ERR_CONTENT_LENGTH_MISMATCH)
                            return StreamingResponse(
                                iter([pdf_content]),
                                media_type="application/pdf",
                                headers={
                                    "Content-Disposition": f'inline; filename="{full_path.stem}.pdf"',
                                    "Content-Length": str(len(pdf_content)),
                                    **FILE_RESPONSE_CORS_HEADERS,
                                },
                            )
                        else:
                            logger.warning(
                                "LibreOffice conversion succeeded but no PDF found"
                            )
                    else:
                        logger.warning(
                            f"LibreOffice conversion failed: {stderr.decode()}"
                        )

            except asyncio.TimeoutError:
                logger.error(f"LibreOffice conversion timed out for PPTX: {full_path}")
            except FileNotFoundError:
                logger.warning(
                    "LibreOffice (soffice) not found, returning original PPTX"
                )
            except Exception as e:
                logger.error(f"Failed to convert PPTX to PDF: {e}")

        # Determine media type
        media_type = "application/octet-stream"
        if filename.endswith((".html", ".htm")):
            media_type = "text/html"
        elif filename.endswith((".css")):
            media_type = "text/css"
        elif filename.endswith((".js")):
            media_type = "application/javascript"
        elif filename.endswith((".jpg", ".jpeg")):
            media_type = "image/jpeg"
        elif filename.endswith((".png")):
            media_type = "image/png"
        elif filename.endswith((".gif")):
            media_type = "image/gif"
        elif filename.endswith((".svg")):
            media_type = "image/svg+xml"
        elif filename.endswith((".webp")):
            media_type = "image/webp"
        elif filename.endswith((".pdf")):
            media_type = "application/pdf"

        return FileResponse(
            path=str(full_path),
            filename=filename,
            media_type=media_type,
            headers=FILE_RESPONSE_CORS_HEADERS,
        )

    except HTTPException:
        # Re-raise HTTP exceptions (like 404 not found)
        raise
    except (PermissionError, OSError) as e:
        # File system error
        logger.error(f"File system error downloading file: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except (ValueError, KeyError, TypeError) as e:
        # Data format error
        logger.error(f"Data format error downloading file: {e}")
        raise HTTPException(status_code=400, detail=f"Data format error: {str(e)}")
    except Exception as e:
        # Other errors, re-raise
        logger.error(f"Unexpected error downloading file: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.get("/preview/{task_id:int}/{file_path:path}", response_model=None)
async def preview_file(
    task_id: int,
    file_path: str,
    db: Session = Depends(get_db),
) -> Any:
    """
    Preview a file with automatic PPTX to HTML conversion.

    This endpoint checks if the file is a PPTX and automatically converts
    it to HTML for browser preview. For other file types, returns the file as-is.
    """
    try:
        logger.info(f"Preview request for task {task_id}, file: {file_path}")

        # Get task from database
        from ..models.task import Task

        task = db.query(Task).filter(Task.id == task_id).first()

        if not task:
            logger.warning(f"Task {task_id} not found for preview")
            raise HTTPException(status_code=404, detail="Task not found")

        # Use the task owner's directory for file access (preferred)
        user_dir = UPLOADS_DIR / f"user_{task.user_id}"

        # Define recursive search function
        def find_file_recursively(directory: Path, filename: str) -> Path | None:
            """Recursively find file in directory"""
            for item in directory.iterdir():
                if item.is_file() and item.name == filename:
                    return item
                elif item.is_dir():
                    result = find_file_recursively(item, filename)
                    if result:
                        return result
            return None

        # Find the file
        found_path: Path | None = None

        # 1) Search in user directory (preferred)
        if user_dir.exists():
            potential_path = user_dir / file_path
            if potential_path.exists():
                found_path = potential_path
            else:
                filename = file_path.split("/")[-1]
                found_path = find_file_recursively(user_dir, filename)

        # 2) Legacy fallback: search in global uploads directory
        if not found_path and UPLOADS_DIR.exists():
            potential_path = UPLOADS_DIR / file_path
            if potential_path.exists():
                found_path = potential_path
            else:
                filename = file_path.split("/")[-1]
                found_path = find_file_recursively(UPLOADS_DIR, filename)

        if not found_path:
            logger.warning(f"File not found for preview: {file_path}")
            raise HTTPException(status_code=404, detail="File not found")

        # Security check: ensure the path is within the task owner's directory or global uploads
        try:
            resolved_path = found_path.resolve()
            resolved_user_dir = user_dir.resolve()
            resolved_global_dir = UPLOADS_DIR.resolve()

            try:
                resolved_path.relative_to(resolved_user_dir)
            except ValueError:
                resolved_path.relative_to(resolved_global_dir)
        except ValueError:
            logger.warning(f"Security check failed for preview path: {file_path}")
            raise HTTPException(status_code=403, detail="Access denied")

        # Get filename early for PPTX processing
        filename = found_path.name

        # Handle PPTX files - convert to PDF using LibreOffice for preview
        if filename.endswith(".pptx"):
            logger.info(f"Converting PPTX to PDF for preview: {found_path}")
            import tempfile

            try:
                # Create a temporary directory for the conversion
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Use LibreOffice to convert PPTX to PDF (async)
                    proc = await asyncio.create_subprocess_exec(
                        "soffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        temp_dir,
                        str(found_path),
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                    )

                    try:
                        stdout, stderr = await asyncio.wait_for(
                            proc.communicate(), timeout=30
                        )
                    except asyncio.TimeoutError:
                        proc.kill()
                        await proc.wait()
                        logger.error(
                            f"LibreOffice conversion timed out for PPTX: {found_path}"
                        )
                        raise

                    if proc.returncode == 0:
                        # Find the generated PDF file
                        pdf_files = list(Path(temp_dir).glob("*.pdf"))
                        if pdf_files:
                            pdf_path = pdf_files[0]
                            logger.info(
                                f"Successfully converted PPTX to PDF: {pdf_path}"
                            )

                            # Read PDF into memory before temp dir is cleaned up
                            pdf_content = pdf_path.read_bytes()

                            # Return streaming response from memory (Content-Length avoids ERR_CONTENT_LENGTH_MISMATCH)
                            return StreamingResponse(
                                iter([pdf_content]),
                                media_type="application/pdf",
                                headers={
                                    "Content-Disposition": f'inline; filename="{found_path.stem}.pdf"',
                                    "Content-Length": str(len(pdf_content)),
                                    **FILE_RESPONSE_CORS_HEADERS,
                                },
                            )
                        else:
                            logger.warning(
                                "LibreOffice conversion succeeded but no PDF found"
                            )
                    else:
                        logger.warning(
                            f"LibreOffice conversion failed: {stderr.decode()}"
                        )

            except asyncio.TimeoutError:
                logger.error(f"LibreOffice conversion timed out for PPTX: {found_path}")
            except FileNotFoundError:
                logger.warning(
                    "LibreOffice (soffice) not found, falling back to text extraction"
                )
            except Exception as e:
                logger.error(f"Failed to convert PPTX to PDF: {e}")

            # Fallback: extract text using python-pptx if LibreOffice fails
            logger.info(f"Falling back to text extraction for PPTX: {found_path}")
            try:
                prs = Presentation(str(found_path))
                html_content = """
                <!DOCTYPE html>
                <html>
                <head>
                    <meta charset="UTF-8">
                    <style>
                        body { font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }
                        h1 { color: #333; border-bottom: 2px solid #4CAF50; padding-bottom: 10px; }
                        h2 { color: #555; margin-top: 30px; }
                        .slide { border: 1px solid #ddd; padding: 20px; margin: 20px 0; background: #f9f9f9; border-radius: 8px; }
                        .slide-number { color: #999; font-size: 12px; margin-top: 10px; }
                        .text-content { white-space: pre-wrap; }
                    </style>
                </head>
                <body>
                    <h1>📊 {filename}</h1>
                """

                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text)

                    if slide_text:
                        html_content += f"""
                        <div class="slide">
                            <h2>Slide {slide_num}</h2>
                            <div class="text-content">
                                {"<br>".join(slide_text)}
                            </div>
                            <div class="slide-number">Slide {slide_num} of {len(prs.slides)}</div>
                        </div>
                        """

                html_content += """
                </body>
                </html>
                """

                return HTMLResponse(content=html_content)

            except Exception as e:
                logger.error(f"Failed to extract text from PPTX: {e}")
                # Fall through to default file response

        # Determine media type
        media_type = "application/octet-stream"
        if filename.endswith((".html", ".htm")):
            media_type = "text/html"
        elif filename.endswith((".css")):
            media_type = "text/css"
        elif filename.endswith((".js")):
            media_type = "application/javascript"
        elif filename.endswith((".jpg", ".jpeg")):
            media_type = "image/jpeg"
        elif filename.endswith((".png")):
            media_type = "image/png"
        elif filename.endswith((".gif")):
            media_type = "image/gif"
        elif filename.endswith((".svg")):
            media_type = "image/svg+xml"
        elif filename.endswith((".webp")):
            media_type = "image/webp"
        elif filename.endswith((".pdf")):
            media_type = "application/pdf"

        headers = {"Content-Disposition": "inline", **FILE_RESPONSE_CORS_HEADERS}
        is_image = filename.lower().endswith(
            (".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg")
        )

        # Images are prone to partial writes / volume sync issues; serve images from memory
        # to avoid ERR_CONTENT_LENGTH_MISMATCH.
        if is_image:
            content = _read_file_with_retry(found_path, max_retries=5, delay=0.1)
            logger.info(
                f"Serving preview image in-memory: {found_path} ({len(content)} bytes)"
            )
            headers["Content-Length"] = str(len(content))
            return Response(content=content, media_type=media_type, headers=headers)

        file_size = found_path.stat().st_size
        if file_size <= PREVIEW_IN_MEMORY_MAX_BYTES:
            content = _read_file_with_retry(found_path, max_retries=5, delay=0.1)
            logger.info(
                f"Serving preview file in-memory: {found_path} ({len(content)} bytes)"
            )
            headers["Content-Length"] = str(len(content))
            return Response(content=content, media_type=media_type, headers=headers)

        logger.info(f"Serving preview file: {found_path}")
        return FileResponse(
            path=str(found_path),
            filename=filename,
            media_type=media_type,
            headers=headers,
        )

    except HTTPException:
        # Re-raise HTTP exceptions
        raise
    except (PermissionError, OSError) as e:
        logger.error(f"File system error in preview: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except Exception as e:
        logger.error(f"Unexpected error in preview: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.get("/preview/{file_path:path}", response_model=None)
async def preview_file_compat(
    file_path: str,
    task_id: int | None = Query(None, description="Task ID when file_path is relative (e.g. output/foo.png)"),
    db: Session = Depends(get_db),
) -> Any:
    """
    Compatibility endpoint.

    Some frontend paths were observed calling:
      /api/files/preview/web_task_47/output/foo.jpg
    or:
      /api/files/preview/output/foo.png?task_id=47

    Extracts task_id from file_path (web_task_XX) or from query, then delegates to the canonical preview endpoint.
    """
    resolved_task_id = _extract_task_id_from_web_task_path(file_path)
    if resolved_task_id is not None:
        return await preview_file(task_id=resolved_task_id, file_path=file_path, db=db)
    if task_id is not None:
        return await preview_file(task_id=task_id, file_path=file_path, db=db)
    raise HTTPException(status_code=400, detail="Invalid file path (missing task id)")


@file_router.get("/public/preview/{task_id:int}/{file_path:path}")
async def public_preview_file(
    task_id: int,
    file_path: str,
    db: Session = Depends(get_db),
) -> FileResponse:
    """
    Public endpoint for previewing task output files.
    This endpoint does not require authentication but validates task ownership.
    """
    try:
        logger.info(f"Public preview request for task {task_id}, file: {file_path}")

        # Get task from database
        from ..models.task import Task

        task = db.query(Task).filter(Task.id == task_id).first()

        if not task:
            logger.warning(f"Task {task_id} not found for public preview")
            raise HTTPException(status_code=404, detail="Task not found")

        # Use the task owner's directory for file access (preferred)
        user_dir = UPLOADS_DIR / f"user_{task.user_id}"

        # Define recursive search function
        def find_file_recursively(directory: Path, filename: str) -> Path | None:
            """Recursively find file in directory"""
            for item in directory.iterdir():
                if item.is_file() and item.name == filename:
                    return item
                elif item.is_dir():
                    result = find_file_recursively(item, filename)
                    if result:
                        return result
            return None

        # Find the file
        found_path: Path | None = None

        # 1) Search in user directory (preferred)
        if user_dir.exists():
            potential_path = user_dir / file_path
            if potential_path.exists():
                found_path = potential_path
            else:
                filename = file_path.split("/")[-1]
                found_path = find_file_recursively(user_dir, filename)

        # 2) Legacy fallback: search in global uploads directory
        if not found_path and UPLOADS_DIR.exists():
            potential_path = UPLOADS_DIR / file_path
            if potential_path.exists():
                found_path = potential_path
            else:
                filename = file_path.split("/")[-1]
                found_path = find_file_recursively(UPLOADS_DIR, filename)

        if not found_path:
            logger.warning(f"File not found for public preview: {file_path}")
            raise HTTPException(status_code=404, detail="File not found")

        # Security check: ensure the path is within the task owner's directory or global uploads
        try:
            resolved_path = found_path.resolve()
            resolved_user_dir = user_dir.resolve()
            resolved_global_dir = UPLOADS_DIR.resolve()

            try:
                resolved_path.relative_to(resolved_user_dir)
            except ValueError:
                resolved_path.relative_to(resolved_global_dir)
        except ValueError:
            logger.warning(
                f"Security check failed for public preview path: {file_path}"
            )
            raise HTTPException(status_code=403, detail="Access denied")

        # Determine media type based on file extension
        filename = found_path.name
        media_type = "application/octet-stream"
        if filename.endswith((".html", ".htm")):
            media_type = "text/html"
        elif filename.endswith((".css")):
            media_type = "text/css"
        elif filename.endswith((".js")):
            media_type = "application/javascript"
        elif filename.endswith((".jpg", ".jpeg")):
            media_type = "image/jpeg"
        elif filename.endswith((".png")):
            media_type = "image/png"
        elif filename.endswith((".gif")):
            media_type = "image/gif"
        elif filename.endswith((".svg")):
            media_type = "image/svg+xml"
        elif filename.endswith((".webp")):
            media_type = "image/webp"

        headers = {"Content-Disposition": "inline", **FILE_RESPONSE_CORS_HEADERS}
        is_image = filename.lower().endswith(
            (".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg")
        )

        # Serve from memory to avoid ERR_CONTENT_LENGTH_MISMATCH (bind mount/NFS or write-in-progress).
        # Images are especially prone to partial writes; always serve images from memory.
        if is_image:
            content = _read_file_with_retry(found_path, max_retries=5, delay=0.1)
            logger.info(
                f"Serving public preview image in-memory: {found_path} ({len(content)} bytes)"
            )
            headers["Content-Length"] = str(len(content))
            return Response(content=content, media_type=media_type, headers=headers)

        # For non-images, serve small files from memory; fall back to FileResponse for large files.
        file_size = found_path.stat().st_size
        if file_size <= PREVIEW_IN_MEMORY_MAX_BYTES:
            content = _read_file_with_retry(found_path, max_retries=5, delay=0.1)
            logger.info(
                f"Serving public preview file in-memory: {found_path} ({len(content)} bytes)"
            )
            headers["Content-Length"] = str(len(content))
            return Response(content=content, media_type=media_type, headers=headers)

        logger.info(f"Serving public preview file: {found_path}")
        return FileResponse(
            path=str(found_path),
            filename=filename,
            media_type=media_type,
            headers=headers,
        )

    except HTTPException:
        # Re-raise HTTP exceptions
        raise
    except (PermissionError, OSError) as e:
        logger.error(f"File system error in public preview: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except Exception as e:
        logger.error(f"Unexpected error in public preview: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.get("/public/preview/{file_path:path}")
async def public_preview_file_compat(
    file_path: str,
    task_id: int | None = Query(None, description="Task ID when file_path is relative (e.g. output/foo.png)"),
    db: Session = Depends(get_db),
) -> Any:
    """
    Compatibility endpoint.

    Some frontend paths were observed calling:
      /api/files/public/preview/web_task_47/output/foo.jpg
    or:
      /api/files/public/preview/output/foo.png?task_id=47

    Extracts task_id from file_path (web_task_XX) or from query, then delegates to the canonical public preview endpoint.
    """
    resolved_task_id = _extract_task_id_from_web_task_path(file_path)
    if resolved_task_id is not None:
        return await public_preview_file(task_id=resolved_task_id, file_path=file_path, db=db)
    if task_id is not None:
        return await public_preview_file(task_id=task_id, file_path=file_path, db=db)
    raise HTTPException(status_code=400, detail="Invalid file path (missing task id)")


@file_router.get("/proxy")
async def proxy_image(url: str) -> Response:
    """
    Proxy remote images to avoid browser CORS / anti-hotlink issues.
    Returns bytes with CORS headers so frontend can load images reliably.
    """
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"}:
        raise HTTPException(status_code=400, detail="Only http/https URLs are supported")
    if _is_private_host(parsed.hostname or ""):
        raise HTTPException(status_code=400, detail="Disallowed host")

    def _fetch() -> tuple[bytes, str]:
        import requests

        headers = {
            "User-Agent": "Mozilla/5.0 (Xagent image proxy)",
            "Accept": "image/*,*/*;q=0.8",
            # Some sites require referer; best-effort use origin.
            "Referer": f"{parsed.scheme}://{parsed.hostname}/",
        }
        # Use a session for better connection handling
        session = requests.Session()
        try:
            with session.get(url, headers=headers, stream=True, timeout=(5, 30), allow_redirects=True) as r:
                r.raise_for_status()
                content_type = r.headers.get("Content-Type", "application/octet-stream")
                expected_length = r.headers.get("Content-Length")

                # Read all content into memory
                content = r.content

                # Verify content length if provided by upstream
                if expected_length and len(content) != int(expected_length):
                    raise IOError(
                        f"Content length mismatch: expected {expected_length}, got {len(content)}"
                    )

                if len(content) > PROXY_IMAGE_MAX_BYTES:
                    raise ValueError("Image too large")

                return content, content_type
        finally:
            session.close()

    try:
        content, content_type = await run_in_threadpool(_fetch)
    except ValueError as e:
        raise HTTPException(status_code=413, detail=str(e))
    except IOError as e:
        logger.warning(f"Proxy fetch IO error: {url} ({e})")
        raise HTTPException(status_code=502, detail=f"Upstream content error: {str(e)}")
    except Exception as e:
        logger.warning(f"Proxy fetch failed: {url} ({e})")
        raise HTTPException(status_code=502, detail="Upstream fetch failed")

    # Prevent intermediaries from compressing/transforming (avoids ERR_CONTENT_LENGTH_MISMATCH)
    headers = {
        **FILE_RESPONSE_CORS_HEADERS,
        "Cache-Control": "public, max-age=3600, no-transform",
        "Content-Encoding": "identity",
        "Content-Length": str(len(content)),
    }
    return Response(content=content, media_type=content_type, headers=headers)


@file_router.delete("/{filename:path}")
async def delete_file(filename: str, user: User = Depends(get_current_user)) -> Dict:
    """Delete uploaded file"""
    try:
        # Get user-specific directory
        user_dir = UPLOADS_DIR / f"user_{user.id}"

        # Define recursive search function once
        def find_file_recursively(directory: Path, filename: str) -> Path | None:
            """Recursively find file in directory"""
            for item in directory.iterdir():
                if item.is_file() and item.name == filename:
                    return item
                elif item.is_dir():
                    result = find_file_recursively(item, filename)
                    if result:
                        return result
            return None

        # Handle both simple filenames and relative paths
        if "/" not in filename:
            # Simple filename - search in user directory first, then fallback to global
            found_path = None

            # Search in user directory first
            if user_dir.exists():
                found_path = find_file_recursively(user_dir, filename)
                if found_path:
                    file_path = found_path
                    logger.info(f"Found file to delete in user directory: {file_path}")

            # If not found in user directory, search in global directory for legacy files
            if not found_path:
                if UPLOADS_DIR.exists():
                    found_path = find_file_recursively(UPLOADS_DIR, filename)
                    if found_path:
                        file_path = found_path
                        logger.info(
                            f"Found file to delete in global directory (legacy): {file_path}"
                        )

            if not found_path:
                logger.warning(f"File to delete not found: {filename}")
                raise HTTPException(status_code=404, detail="File not found")
        else:
            # Relative path - try direct path first, then fallback to recursive search
            file_path = user_dir / filename

            if not file_path.exists():
                # Fallback: search recursively for the filename with path preference
                target_filename = filename.split("/")[-1]

                # Enhanced search that prefers matching path structure (same logic as download)
                def find_file_with_path_preference(
                    directory: Path, filename: str, original_path: str
                ) -> Path | None:
                    """Find file recursively, preferring paths that match the original structure"""
                    best_match = None
                    best_score = 0

                    def score_match(candidate_path: Path, target_path: str) -> int:
                        """Score how well a candidate path matches the target structure"""
                        # Get relative path from user directory
                        try:
                            relative_path = candidate_path.relative_to(user_dir)
                        except ValueError:
                            # Fallback to using absolute path parts
                            relative_path = candidate_path

                        candidate_parts = relative_path.parts
                        target_parts = target_path.split("/")

                        score = 0
                        # Check if path parts match
                        for i, target_part in enumerate(
                            target_parts[:-1]
                        ):  # Exclude filename
                            if (
                                i < len(candidate_parts) - 1
                            ):  # Exclude filename from candidate
                                candidate_part = candidate_parts[i]
                                if target_part == candidate_part:
                                    score += 10
                                elif (
                                    target_part.replace("web_task_", "task_")
                                    == candidate_part
                                ):
                                    score += (
                                        5  # Partial match for task prefix differences
                                    )
                                elif (
                                    candidate_part.replace("web_task_", "task_")
                                    == target_part
                                ):
                                    score += 5  # Reverse match
                        return score

                    for item in directory.iterdir():
                        if item.is_file() and item.name == filename:
                            # Score this match
                            current_score = score_match(item, original_path)
                            if current_score > best_score:
                                best_match = item
                                best_score = current_score
                                logger.debug(
                                    f"Found better match for delete with score {current_score}: {item}"
                                )
                        elif item.is_dir():
                            result = find_file_with_path_preference(
                                item, filename, original_path
                            )
                            if result:
                                current_score = score_match(result, original_path)
                                if current_score > best_score:
                                    best_match = result
                                    best_score = current_score

                    return best_match

                found_path = find_file_with_path_preference(
                    user_dir, target_filename, filename
                )
                if found_path:
                    file_path = found_path
                    logger.info(
                        f"Found file to delete by enhanced recursive search: {file_path}"
                    )
                else:
                    logger.warning(f"File to delete not found: {filename}")
                    raise HTTPException(status_code=404, detail="File not found")
            else:
                logger.info(f"Found file to delete by direct path: {file_path}")

        # Security check: ensure the path is within allowed directories
        try:
            resolved_path = file_path.resolve()
            resolved_user_dir = user_dir.resolve()
            resolved_global_dir = UPLOADS_DIR.resolve()

            # Allow both user directory and global directory for legacy files
            try:
                resolved_path.relative_to(resolved_user_dir)
            except ValueError:
                try:
                    resolved_path.relative_to(resolved_global_dir)
                except ValueError:
                    logger.warning(f"Security check failed for delete path: {filename}")
                    raise HTTPException(status_code=403, detail="Access denied")
        except ValueError:
            logger.warning(f"Security check failed for delete path: {filename}")
            raise HTTPException(status_code=403, detail="Access denied")

        if not file_path.exists():
            logger.warning(f"File to delete not found: {file_path}")
            raise HTTPException(status_code=404, detail="File not found")

        file_path.unlink()
        logger.info(f"Deleted file: {file_path}")

        return {"success": True, "message": f"File {filename} deleted successfully"}

    except HTTPException:
        # Re-raise HTTP exceptions (like 404 not found)
        raise
    except (PermissionError, OSError) as e:
        # File system permission error
        logger.error(f"File system error deleting file: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except (ValueError, KeyError, TypeError) as e:
        # Data format error
        logger.error(f"Data format error deleting file: {e}")
        raise HTTPException(status_code=400, detail=f"Data format error: {str(e)}")
    except Exception as e:
        # Other errors, re-raise
        logger.error(f"Unexpected error deleting file: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")
